#!/usr/bin/env python3
"""Convert a PDF into a PPTX, one slide per page.

Two engines, picked with --engine:

* ``image`` (default) — render each PDF page to a raster image (PyMuPDF)
  and place it full-bleed on a blank slide. Pixel-perfect appearance —
  fonts, LaTeX equations, vector figures — at the cost of editable text.
  Requires: PyMuPDF (fitz), python-pptx (both on the cluster).

* ``libreoffice`` — drive headless LibreOffice to import the PDF as an
  Impress document and save it as PPTX. Yields *editable* text + embedded
  figures, but text is fragmented per-span and Beamer math degrades to
  glyph-level text (no native PowerPoint equations). Requires a `soffice`
  binary on PATH (`sudo dnf install libreoffice-impress libreoffice-draw`).

There is no free, fully-clean editable conversion: PDF carries no slide
structure, so reconstruction is always heuristic.

Usage:
    python scripts/pdf_to_pptx.py input.pdf                      # image engine
    python scripts/pdf_to_pptx.py input.pdf --engine libreoffice # editable
    python scripts/pdf_to_pptx.py input.pdf -o out.pptx --dpi 200
    python scripts/pdf_to_pptx.py docs/*.pdf --engine libreoffice # batch
"""
import argparse
import io
import shutil
import subprocess
import sys
from pathlib import Path

from pptx import Presentation
from pptx.util import Emu

# PowerPoint measures in EMU; 914400 EMU = 1 inch. PDF points are 1/72 inch.
EMU_PER_INCH = 914400
EMU_PER_PT = EMU_PER_INCH // 72


def convert_image(pdf_path: Path, pptx_path: Path, dpi: int) -> int:
    """Render every page of *pdf_path* onto slides and save to *pptx_path*.

    Returns the number of slides written.
    """
    import fitz  # PyMuPDF — imported lazily so the libreoffice engine needn't have it

    doc = fitz.open(pdf_path)
    if doc.page_count == 0:
        raise ValueError(f"{pdf_path} has no pages")

    prs = Presentation()
    # Size the deck to the first page; mixed-size PDFs still render correctly
    # because each image is fit to the slide, but uniform pages look cleanest.
    first = doc.load_page(0)
    prs.slide_width = Emu(round(first.rect.width * EMU_PER_PT))
    prs.slide_height = Emu(round(first.rect.height * EMU_PER_PT))
    blank_layout = prs.slide_layouts[6]  # fully blank

    zoom = dpi / 72.0
    matrix = fitz.Matrix(zoom, zoom)

    for i in range(doc.page_count):
        page = doc.load_page(i)
        pix = page.get_pixmap(matrix=matrix, alpha=False)
        img = io.BytesIO(pix.tobytes("png"))

        slide = prs.slides.add_slide(blank_layout)
        # Fit image to the page's own size (handles mixed page sizes), centered.
        pw = round(page.rect.width * EMU_PER_PT)
        ph = round(page.rect.height * EMU_PER_PT)
        left = (prs.slide_width - pw) // 2
        top = (prs.slide_height - ph) // 2
        slide.shapes.add_picture(img, Emu(left), Emu(top), width=Emu(pw), height=Emu(ph))

    pptx_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(pptx_path)
    doc.close()
    return prs.slides and len(prs.slides._sldIdLst) or 0


def convert_libreoffice(pdf_path: Path, pptx_path: Path, profile: Path) -> int:
    """Convert via headless LibreOffice into an *editable* PPTX.

    Returns the number of slides written. *profile* is an isolated
    UserInstallation dir so a running LibreOffice / sandbox can't block us.
    """
    soffice = shutil.which("soffice") or shutil.which("libreoffice")
    if soffice is None:
        raise RuntimeError(
            "soffice not found — install with "
            "`sudo dnf install libreoffice-impress libreoffice-draw`")

    pptx_path.parent.mkdir(parents=True, exist_ok=True)
    # LibreOffice writes <stem>.pptx into --outdir; it ignores any other name.
    # NOTE: passing an explicit output filter (pptx:"impress_MS_PowerPoint...")
    # triggers a store error (0x81a); use the bare `pptx` target. The
    # --infilter forces import as Impress (editable) rather than Draw.
    cmd = [
        soffice, "--headless", "--norestore",
        f"-env:UserInstallation=file://{profile}",
        "--convert-to", "pptx",
        "--infilter=impress_pdf_import",
        str(pdf_path),
        "--outdir", str(pptx_path.parent),
    ]
    proc = subprocess.run(cmd, capture_output=True, text=True, timeout=600)
    produced = pptx_path.parent / (pdf_path.stem + ".pptx")
    if not produced.is_file():
        raise RuntimeError(
            "LibreOffice produced no output. "
            f"stdout={proc.stdout.strip()!r} stderr={proc.stderr.strip()!r}")
    if produced != pptx_path:
        produced.replace(pptx_path)

    try:
        return len(Presentation(pptx_path).slides._sldIdLst)
    except Exception:  # noqa: BLE001 — slide count is cosmetic; file already exists
        return -1


def main(argv=None):
    ap = argparse.ArgumentParser(description="Convert PDF(s) to PPTX, one slide per page.")
    ap.add_argument("pdf", nargs="+", type=Path, help="input PDF file(s)")
    ap.add_argument("-o", "--out", type=Path, default=None,
                    help="output .pptx (single input only; default: alongside the PDF)")
    ap.add_argument("--engine", choices=("image", "libreoffice"), default="image",
                    help="image: pixel-perfect, non-editable (default). "
                         "libreoffice: editable text, fragmented layout.")
    ap.add_argument("--dpi", type=int, default=200,
                    help="[image engine] render resolution (default 200; 150 smaller, 300 sharper)")
    ap.add_argument("--profile", type=Path, default=Path("/tmp/lo_profile_pdf2pptx"),
                    help="[libreoffice engine] isolated LibreOffice user profile dir")
    args = ap.parse_args(argv)

    if args.out is not None and len(args.pdf) > 1:
        ap.error("-o/--out cannot be used with multiple input PDFs")

    rc = 0
    for pdf in args.pdf:
        if not pdf.is_file():
            print(f"skip: {pdf} is not a file", file=sys.stderr)
            rc = 1
            continue
        out = args.out if args.out is not None else pdf.with_suffix(".pptx")
        try:
            if args.engine == "image":
                n = convert_image(pdf, out, args.dpi)
                detail = f"{n} slides @ {args.dpi} dpi, image"
            else:
                n = convert_libreoffice(pdf, out, args.profile)
                detail = f"{n} slides, editable (libreoffice)"
        except Exception as e:  # noqa: BLE001 — surface any conversion failure per-file
            print(f"FAIL: {pdf} -> {e}", file=sys.stderr)
            rc = 1
            continue
        print(f"OK: {pdf} -> {out} ({detail})")
    return rc


if __name__ == "__main__":
    raise SystemExit(main())

#!/usr/bin/env python3
"""Add MICOM explanation slides to progress_report_hamilton_kegg.pptx."""
import sys as _sys, pathlib as _pathlib  # noqa: E402  [nife-pathshim]
_sys.path.insert(0, str(_pathlib.Path(__file__).resolve().parents[2]))  # repo root: bare sibling imports
import argparse
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.oxml.ns import qn
from copy import deepcopy
import lxml.etree as etree
from pathlib import Path
import shutil, datetime

_ap = argparse.ArgumentParser()
_ap.add_argument('--replace', action='store_true',
                 help='Replace last 3 slides (MICOM slides) instead of appending')
_ap.add_argument('--n-replace', type=int, default=3)
_args = _ap.parse_args()

PPTX = Path('/home/nishioka/IKM_Hiwi/nife/results/dieckow_cr/progress_report_hamilton_kegg.pptx')

# ── Color palette ──────────────────────────────────────────────────────────────
C_BLUE     = RGBColor(0x1a, 0x3d, 0x7a)   # dark blue  (header bg)
C_LBLUE    = RGBColor(0x3a, 0x6e, 0xc8)   # mid blue
C_GREEN    = RGBColor(0x1a, 0x73, 0x40)   # dark green (good)
C_RED      = RGBColor(0xcc, 0x22, 0x22)   # red        (bad)
C_ORANGE   = RGBColor(0xe0, 0x70, 0x20)   # orange     (warn)
C_WHITE    = RGBColor(0xff, 0xff, 0xff)
C_BLACK    = RGBColor(0x11, 0x11, 0x11)
C_LGRAY    = RGBColor(0xf0, 0xf2, 0xf5)
C_DGRAY    = RGBColor(0x55, 0x55, 0x55)
C_TEAL     = RGBColor(0x00, 0x7a, 0x8a)

W = Inches(13.33)
H = Inches(7.5)

# ── Helpers ────────────────────────────────────────────────────────────────────

def add_rect(slide, left, top, width, height, fill=None, line=None, line_w=None):
    shape = slide.shapes.add_shape(1, left, top, width, height)  # MSO_SHAPE_TYPE.RECTANGLE
    fill_fmt = shape.fill
    if fill:
        fill_fmt.solid()
        fill_fmt.fore_color.rgb = fill
    else:
        fill_fmt.background()
    line_fmt = shape.line
    if line:
        line_fmt.color.rgb = line
        if line_w:
            line_fmt.width = line_w
    else:
        line_fmt.fill.background()
    shape.shadow.inherit = False
    return shape


def add_textbox(slide, text, left, top, width, height,
                font_size=18, bold=False, color=C_BLACK,
                align=PP_ALIGN.LEFT, wrap=True, italic=False):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return txBox


def add_header(slide, title, subtitle=None):
    """Blue header bar with white title."""
    add_rect(slide, 0, 0, W, Inches(1.1), fill=C_BLUE)
    add_textbox(slide, title, Inches(0.4), Inches(0.12), Inches(12.5), Inches(0.6),
                font_size=28, bold=True, color=C_WHITE, align=PP_ALIGN.LEFT)
    if subtitle:
        add_textbox(slide, subtitle, Inches(0.4), Inches(0.72), Inches(12.5), Inches(0.35),
                    font_size=14, color=RGBColor(0xcc, 0xdd, 0xff), align=PP_ALIGN.LEFT)


def add_bullet_box(slide, title, bullets, left, top, width, height,
                   title_color=C_BLUE, bullet_size=15, title_size=17):
    """Rounded box with title + bullet list."""
    add_rect(slide, left, top, width, Inches(0.42), fill=title_color)
    add_textbox(slide, title, left + Inches(0.15), top + Inches(0.05),
                width - Inches(0.3), Inches(0.32),
                font_size=title_size, bold=True, color=C_WHITE)
    add_rect(slide, left, top + Inches(0.42), width, height - Inches(0.42),
             fill=C_LGRAY, line=title_color, line_w=Pt(1))
    # Bullets as multiline textbox
    y = top + Inches(0.52)
    for bullet in bullets:
        indent = bullet.startswith('  ')
        text = ('    ' if indent else '• ') + bullet.lstrip()
        sz = bullet_size - 2 if indent else bullet_size
        clr = C_DGRAY if indent else C_BLACK
        add_textbox(slide, text, left + Inches(0.1), y,
                    width - Inches(0.2), Inches(0.4),
                    font_size=sz, color=clr)
        y += Inches(0.33 if indent else 0.38)


def add_table_row(slide, cells, left, top, width, height,
                  bg=C_WHITE, text_color=C_BLACK, font_size=13, bold=False):
    col_w = width / len(cells)
    for k, cell in enumerate(cells):
        add_rect(slide, left + k * col_w, top, col_w, height,
                 fill=bg, line=RGBColor(0xcc, 0xcc, 0xcc), line_w=Pt(0.5))
        add_textbox(slide, cell, left + k * col_w + Inches(0.08), top + Inches(0.04),
                    col_w - Inches(0.16), height - Inches(0.08),
                    font_size=font_size, bold=bold, color=text_color, align=PP_ALIGN.CENTER)


prs = Presentation(PPTX)

# Remove last N slides if --replace
if _args.replace:
    n = _args.n_replace
    total = len(prs.slides)
    sldIdLst = prs.slides._sldIdLst
    for _ in range(min(n, total)):
        idx = len(sldIdLst) - 1
        rId = sldIdLst[idx].get(qn('r:id'))
        prs.part.drop_rel(rId)
        del sldIdLst[idx]
    print(f'Removed last {min(n, total)} slides (was {total}, now {len(prs.slides)})')

blank_layout = prs.slide_layouts[6]   # Blank

# ══════════════════════════════════════════════════════════════════════════════
# Slide A: なぜ MICOM を使ったか — 失敗した手法と問題点
# ══════════════════════════════════════════════════════════════════════════════
slideA = prs.slides.add_slide(blank_layout)
add_rect(slideA, 0, 0, W, H, fill=C_WHITE)
add_header(slideA, 'Sign Prior Layer 3: なぜ MICOM を選んだか',
           subtitle='— 単種 pFBA / MacArthur / Growth-Rate Suppression の失敗と MICOM への移行 —')

# 左: 失敗手法
add_bullet_box(slideA, '❌  MacArthur コサイン類似度',
    ['competition[i,j] = (c_i · c_j) / (|c_i| |c_j|)',
     '失敗理由:',
     '  口腔菌はすべて generalist（アミノ酸・糖を全員が使う）',
     '  → cosine ≈ 1.0 for all pairs → 全ペアが競合と誤判定',
     '  MacArthur (1970) の前提「種ごとに異なるニッチ資源」が不成立'],
    Inches(0.3), Inches(1.25), Inches(5.9), Inches(2.7),
    title_color=C_RED)

add_bullet_box(slideA, '❌  Growth Rate Suppression (GRS)',
    ['competition[i,j] = (μ_i(full) − μ_i(depleted)) / μ_i(full)',
     '失敗理由:',
     '  v2 唾液培地が極度に乏しく、菌は barely surviving',
     '  → j の uptake で何かが枯渇すると i が完全に成長停止',
     '  81 / 90 ペアで Δμ/μ ≈ 100%  → 非特異的すぎて使えない'],
    Inches(0.3), Inches(4.15), Inches(5.9), Inches(2.8),
    title_color=C_RED)

# 右: 単種 pFBA の限界
add_bullet_box(slideA, '⚠️  単種 pFBA（v1 / v2）の限界',
    ['各菌を独立に FBA → 代謝物の「移動可能性」しか見ない',
     '実際に j から i へ代謝物が流れているか不明',
     'v1: 培地が現実の 100 倍豊富 → シグナルは強いが非現実的',
     'v2: 培地が乏しすぎ → cross-feeding シグナルが弱い',
     '競合シグナル: MacArthur/GRS ともに失敗'],
    Inches(6.5), Inches(1.25), Inches(6.5), Inches(2.7),
    title_color=C_ORANGE)

# 右下: MICOM の結論
add_rect(slideA, Inches(6.5), Inches(4.15), Inches(6.5), Inches(2.8),
         fill=RGBColor(0xe8, 0xf4, 0xec), line=C_GREEN, line_w=Pt(2))
add_textbox(slideA, '✅  MICOM の採用理由',
            Inches(6.65), Inches(4.25), Inches(6.2), Inches(0.4),
            font_size=17, bold=True, color=C_GREEN)
for y_off, line in enumerate([
    '10 菌を同時に入れた community FBA を解く',
    '→ 実際に代謝物が j から i へ「流れているか」を直接確認',
    '→ 口腔菌 generalist 問題を回避（ニッチ競合でなくフラックス）',
    '→ Streptococcus → Veillonella 乳酸 cross-feeding を直接実証',
    '→ Sign Agreement: 100% (36/36)  vs  v1: 88% (29/33)',
]):
    clr = C_GREEN if '100%' in line or '直接実証' in line else C_BLACK
    bld = '100%' in line
    add_textbox(slideA, '• ' + line, Inches(6.65), Inches(4.75) + Inches(0.37 * y_off),
                Inches(6.2), Inches(0.35), font_size=14, color=clr, bold=bld)

# ══════════════════════════════════════════════════════════════════════════════
# Slide B: MICOM の仕組み — Community FBA + Cooperative Tradeoff
# ══════════════════════════════════════════════════════════════════════════════
slideB = prs.slides.add_slide(blank_layout)
add_rect(slideB, 0, 0, W, H, fill=C_WHITE)
add_header(slideB, 'MICOM のメカニズム: Community FBA + Cooperative Tradeoff',
           subtitle='Diener et al. 2020, mSystems 5:e00606-19')

# Step boxes
steps = [
    ('Step 1', '10 guild モデルを統合', 'AGORA2 SBML × 10 guild\n→ 1 つの community モデル\n（共有培地プール）'),
    ('Step 2', 'Cooperative Tradeoff', 'maximize  min_i(μ_i / μ_i_max)\ns.t.  Σμ_i ≥ τ × Σμ_i_max\nτ = 0.5 (Diener 2020 default)'),
    ('Step 3', 'Flux 方向を抽出', '各 exchange rxn につき\n  分泌菌 (flux > 0)\n  消費菌 (flux < 0)\n→ 実際の代謝物移動'),
    ('Step 4', 'Sign Prior 生成', 'pos[i,j] += min(src, |tgt|)\n（flux 量で重み付け）\nneg[i,j]: H₂O₂, H₂S → 全菌'),
]
for k, (num, title, body) in enumerate(steps):
    lft = Inches(0.3 + k * 3.3)
    add_rect(slideB, lft, Inches(1.3), Inches(3.1), Inches(0.45), fill=C_LBLUE)
    add_textbox(slideB, f'{num}: {title}', lft + Inches(0.1), Inches(1.35),
                Inches(2.9), Inches(0.35), font_size=15, bold=True, color=C_WHITE)
    add_rect(slideB, lft, Inches(1.75), Inches(3.1), Inches(1.9),
             fill=C_LGRAY, line=C_LBLUE, line_w=Pt(1))
    add_textbox(slideB, body, lft + Inches(0.1), Inches(1.82),
                Inches(2.9), Inches(1.75), font_size=13, color=C_BLACK)
    if k < 3:
        add_textbox(slideB, '→', lft + Inches(3.15), Inches(2.35),
                    Inches(0.2), Inches(0.4), font_size=20, bold=True, color=C_LBLUE,
                    align=PP_ALIGN.CENTER)

# Key example: Lactate cross-feeding
add_rect(slideB, Inches(0.3), Inches(3.85), Inches(8.0), Inches(2.6),
         fill=RGBColor(0xe8, 0xf4, 0xec), line=C_GREEN, line_w=Pt(1.5))
add_textbox(slideB, '実測 Community Flux の例: Streptococcus → Veillonella 乳酸 cross-feeding',
            Inches(0.45), Inches(3.9), Inches(7.7), Inches(0.4),
            font_size=15, bold=True, color=C_GREEN)
code = ('EX_lac_L(e):\n'
        '  Bacilli   (Streptococcus)  secretion  = +97.7 mmol/gDW/h\n'
        '  Negativicutes (Veillonella) uptake     = −97.9 mmol/gDW/h\n'
        '  → pos[Negativicutes, Bacilli] += 97.7   (高重み)')
add_textbox(slideB, code, Inches(0.45), Inches(4.38), Inches(7.7), Inches(1.9),
            font_size=13, color=C_BLACK, italic=False)

# Why cooperative tradeoff works for generalists
add_rect(slideB, Inches(8.6), Inches(3.85), Inches(4.4), Inches(2.6),
         fill=RGBColor(0xf0, 0xf0, 0xff), line=C_LBLUE, line_w=Pt(1.5))
add_textbox(slideB, 'なぜ Generalist 問題を解決？',
            Inches(8.75), Inches(3.9), Inches(4.1), Inches(0.4),
            font_size=15, bold=True, color=C_BLUE)
for y, txt in enumerate([
    '全員が同じ資源を使えても',
    '→ community FBA で資源を「分配」',
    '→ 各菌が実際に使う経路が決まる',
    '',
    'MacArthur: 使える？ (static)',
    'MICOM: 実際に使った？ (dynamic)',
]):
    add_textbox(slideB, txt, Inches(8.75), Inches(4.38) + Inches(0.35 * y),
                Inches(4.1), Inches(0.32), font_size=13,
                color=C_BLUE if 'MICOM' in txt else C_BLACK,
                bold='MICOM' in txt)

# ══════════════════════════════════════════════════════════════════════════════
# Slide C: 結果比較 + Perfect 版の改良点
# ══════════════════════════════════════════════════════════════════════════════
slideC = prs.slides.add_slide(blank_layout)
add_rect(slideC, 0, 0, W, H, fill=C_WHITE)
add_header(slideC, 'MICOM Sign Prior: 結果比較と "Perfect" 版の改良',
           subtitle='sign agreement / constrained pairs / LOO-RMSE')

# Table: sign agreement comparison
headers = ['手法', 'Sign Agreement', '制約ペア数', '備考']
add_table_row(slideC, headers, Inches(0.3), Inches(1.3), Inches(7.5), Inches(0.45),
              bg=C_BLUE, text_color=C_WHITE, font_size=13, bold=True)
rows = [
    ('文献 L1+L2 のみ',      '45%  (5/11)',   '11 / 45', '文献ベース sign のみ'),
    ('単種 pFBA v1',         '88%  (29/33)',  '33 / 45', 'Blood-plasma medium'),
    ('単種 pFBA v2',         '81%  (21/26)',  '26 / 45', '唾液相当(乏しすぎ)'),
    ('MICOM (τ=0.5)',       '100% (72/72)',  '72 / 90', '★ Flux-weighted, toxin→全菌'),
]
row_colors = [C_LGRAY, C_WHITE, C_LGRAY,
              RGBColor(0xe8, 0xf4, 0xec)]
text_colors = [C_BLACK] * 3 + [C_GREEN]
for k, (row, bg, tc) in enumerate(zip(rows, row_colors, text_colors)):
    bld = k == 4
    add_table_row(slideC, list(row), Inches(0.3), Inches(1.75) + Inches(0.42 * k),
                  Inches(7.5), Inches(0.42), bg=bg, text_color=tc,
                  font_size=13, bold=bld)

# Perfect version changes table
add_textbox(slideC, '"Perfect" 版の 3 つの改良点',
            Inches(0.3), Inches(4.15), Inches(7.5), Inches(0.4),
            font_size=16, bold=True, color=C_BLUE)
p_headers = ['改良点', 'v0 (初期)', 'Perfect 版', '理由']
add_table_row(slideC, p_headers, Inches(0.3), Inches(4.6), Inches(7.5), Inches(0.42),
              bg=C_BLUE, text_color=C_WHITE, font_size=12, bold=True)
p_rows = [
    ('Tradeoff τ',     'τ = 0.3',       'τ = 0.5',          'Diener 2020 推奨値'),
    ('Cross-feeding 重み', 'バイナリ +1', 'min(src, |tgt|)',  '高 flux を正当評価'),
    ('毒素シグナル', 'consumer のみ', '全 guild に伝播',    '毒素は拡散・代謝しない'),
]
for k, row in enumerate(p_rows):
    bg = C_LGRAY if k % 2 == 0 else C_WHITE
    add_table_row(slideC, list(row), Inches(0.3), Inches(5.02) + Inches(0.42 * k),
                  Inches(7.5), Inches(0.42), bg=bg, font_size=12)

# Right: LOO-CV status
add_rect(slideC, Inches(8.1), Inches(1.3), Inches(4.9), Inches(5.65),
         fill=RGBColor(0xf5, 0xf8, 0xff), line=C_BLUE, line_w=Pt(1.5))
add_textbox(slideC, 'LOO-CV 最終結果 (2026-05-20)',
            Inches(8.25), Inches(1.38), Inches(4.6), Inches(0.4),
            font_size=15, bold=True, color=C_BLUE)
loo_results = [
    ('Hamilton MICOM τ=0.5', '0.0502', '72/72', C_GREEN),
    ('Hamilton AGORA-v1',    '0.0562', '66/70', C_BLACK),
    ('Hamilton no-prior',    '0.0595', '20/28', C_DGRAY),
    ('gLV MICOM τ=0.5',     '0.0513', '71/72', C_BLACK),
    ('gLV AGORA-v1',        '0.0499', '52/53', C_BLACK),
    ('gLV α=0.25 ★best',   '0.0490', '32/33', C_GREEN),
]
add_table_row(slideC, ['モデル', 'RMSE', 'SA'],
              Inches(8.1), Inches(1.85), Inches(4.9), Inches(0.38),
              bg=C_BLUE, text_color=C_WHITE, font_size=12, bold=True)
for k, (name, rmse, sa, clr) in enumerate(loo_results):
    bg = RGBColor(0xe8, 0xf4, 0xec) if clr == C_GREEN else (C_LGRAY if k % 2 == 0 else C_WHITE)
    add_table_row(slideC, [name, rmse, sa],
                  Inches(8.1), Inches(2.23) + Inches(0.38 * k), Inches(4.9), Inches(0.38),
                  bg=bg, text_color=clr, font_size=12, bold=(clr == C_GREEN))

# ── Save ────────────────────────────────────────────────────────────────────
out = PPTX
shutil.copy(PPTX, PPTX.with_suffix('.bak.pptx'))
prs.save(out)
print(f'Saved {out}  ({len(prs.slides)} slides total)')
print('Backup: ' + str(PPTX.with_suffix('.bak.pptx')))

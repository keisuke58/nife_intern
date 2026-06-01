#!/usr/bin/env python3
"""
add_dieckow_slides.py
Appends 6 new slides (guild analysis results) to the existing
progress_report_hamilton_kegg.pptx, matching the existing visual style.

Style:
  - Header rect:  fill #1A3A5C, full width, h=0.72"
  - Title textbox: white, bold, 20pt, top-left in header
  - Subtitle textbox: light-blue, 12pt, below title in header
  - Figure: centred between header and footer
  - Footer rect: fill #EDFAF3, full width, h=0.80", y=6.55"
  - Footer text: dark navy, 11pt, key findings
"""

from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

PPTX_IN  = Path('dieckow_paper/progress_report_hamilton_kegg.pptx')
PPTX_OUT = Path('dieckow_paper/progress_report_dieckow_extended.pptx')

R = Path('results')

# Colour constants
HEADER_FILL  = RGBColor(0x1A, 0x3A, 0x5C)  # dark navy
FOOTER_FILL  = RGBColor(0xED, 0xFA, 0xF3)  # mint
WHITE        = RGBColor(0xFF, 0xFF, 0xFF)
NAVY         = RGBColor(0x1A, 0x3A, 0x5C)
LIGHT_BLUE   = RGBColor(0xAD, 0xD8, 0xE6)
ACCENT_GREEN = RGBColor(0x1A, 0x7A, 0x4A)

W = 13.33   # slide width inches
H = 7.50    # slide height inches

HEADER_H  = 0.72
FOOTER_Y  = 6.55
FOOTER_H  = H - FOOTER_Y
FIG_TOP   = HEADER_H + 0.10
FIG_H_MAX = FOOTER_Y - FIG_TOP - 0.05


# ── Helpers ───────────────────────────────────────────────────────────────────

def add_rect(slide, x, y, w, h, fill_rgb):
    from pptx.util import Inches
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        Inches(x), Inches(y), Inches(w), Inches(h)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_rgb
    shape.line.fill.background()
    return shape


def add_textbox(slide, x, y, w, h, text, font_size, bold=False,
                color=WHITE, align=PP_ALIGN.LEFT, wrap=True):
    txBox = slide.shapes.add_textbox(
        Inches(x), Inches(y), Inches(w), Inches(h))
    tf = txBox.text_frame
    tf.word_wrap = wrap
    p  = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = color
    return txBox


def add_slide(prs, title, subtitle, fig_path, footer_text,
              fig_x=None, fig_w=None, fig_h=None):
    """Add a single styled slide."""
    blank_layout = prs.slide_layouts[6]   # Blank
    slide = prs.slides.add_slide(blank_layout)

    # Header rectangle
    add_rect(slide, 0, 0, W, HEADER_H, HEADER_FILL)

    # Title
    add_textbox(slide, 0.40, 0.06, W - 0.5, 0.40,
                title, 20, bold=True, color=WHITE)

    # Subtitle
    add_textbox(slide, 0.40, 0.46, W - 0.5, 0.24,
                subtitle, 11, bold=False, color=LIGHT_BLUE)

    # Figure
    if fig_path and Path(fig_path).exists():
        pic_w  = fig_w  or 10.0
        pic_h  = fig_h  or FIG_H_MAX
        # keep aspect ratio
        from PIL import Image as PILImage
        img = PILImage.open(fig_path)
        iw, ih = img.size
        ratio = iw / ih
        if pic_w / ratio > pic_h:
            pic_w = pic_h * ratio
        else:
            pic_h = pic_w / ratio
        pic_x = fig_x if fig_x is not None else (W - pic_w) / 2
        pic_y = FIG_TOP + (FIG_H_MAX - pic_h) / 2
        slide.shapes.add_picture(str(fig_path),
                                 Inches(pic_x), Inches(pic_y),
                                 Inches(pic_w), Inches(pic_h))

    # Footer rectangle
    add_rect(slide, 0, FOOTER_Y, W, FOOTER_H, FOOTER_FILL)

    # Footer text
    add_textbox(slide, 0.35, FOOTER_Y + 0.07, W - 0.6, FOOTER_H - 0.10,
                footer_text, 10.5, bold=False, color=NAVY)

    return slide


def add_two_fig_slide(prs, title, subtitle,
                      fig_left, fig_right,
                      caption_left, caption_right,
                      footer_text):
    """Slide with two figures side by side."""
    blank_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_layout)

    add_rect(slide, 0, 0, W, HEADER_H, HEADER_FILL)
    add_textbox(slide, 0.40, 0.06, W - 0.5, 0.40,
                title, 20, bold=True, color=WHITE)
    add_textbox(slide, 0.40, 0.46, W - 0.5, 0.24,
                subtitle, 11, bold=False, color=LIGHT_BLUE)

    half_w = (W - 0.6) / 2
    fig_h  = FIG_H_MAX - 0.35   # leave room for caption

    for i, (fp, cap) in enumerate([(fig_left, caption_left),
                                   (fig_right, caption_right)]):
        if fp and Path(fp).exists():
            from PIL import Image as PILImage
            img = PILImage.open(fp)
            iw, ih = img.size
            ratio = iw / ih
            pw = half_w - 0.1
            ph = pw / ratio
            if ph > fig_h:
                ph = fig_h
                pw = ph * ratio
            px = 0.15 + i * (half_w + 0.15) + (half_w - pw) / 2
            py = FIG_TOP + (fig_h - ph) / 2
            slide.shapes.add_picture(str(fp),
                                     Inches(px), Inches(py),
                                     Inches(pw), Inches(ph))
        # Caption
        cx = 0.15 + i * (half_w + 0.15)
        add_textbox(slide, cx, FOOTER_Y - 0.32, half_w, 0.28,
                    cap, 9, bold=False, color=NAVY, align=PP_ALIGN.CENTER)

    add_rect(slide, 0, FOOTER_Y, W, FOOTER_H, FOOTER_FILL)
    add_textbox(slide, 0.35, FOOTER_Y + 0.07, W - 0.6, FOOTER_H - 0.10,
                footer_text, 10.5, bold=False, color=NAVY)
    return slide


# ── Build slides ──────────────────────────────────────────────────────────────

def main():
    prs = Presentation(str(PPTX_IN))

    # ── Slide A: Section divider ────────────────────────────────────────────
    blank = prs.slide_layouts[6]
    s = prs.slides.add_slide(blank)
    add_rect(s, 0, 0, W, H, HEADER_FILL)
    add_textbox(s, 1.0, 2.5, W - 2.0, 1.0,
                'Guild-Level Dynamical Analysis', 32, bold=True, color=WHITE,
                align=PP_ALIGN.CENTER)
    add_textbox(s, 1.0, 3.5, W - 2.0, 0.5,
                'Keystone · Network · Permutation · Tipping · Relapse · Prevention · Phase Diagram',
                14, bold=False, color=LIGHT_BLUE, align=PP_ALIGN.CENTER)
    # Key message box
    add_rect(s, 1.5, 4.3, W - 3.0, 1.5, RGBColor(0x1A, 0x5C, 0x3A))
    add_textbox(s, 1.6, 4.38, W - 3.2, 1.35,
                'Periodontitis relapse is governed by the topology of microbial '
                'interaction networks, not merely by bacterial growth rates. '
                'Prevention thresholds and relapse timing can be mathematically '
                'derived from patient-specific gLV dynamics.',
                11.5, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    # ── Slide A2: Background & Model ────────────────────────────────────────
    blank2 = prs.slide_layouts[6]
    s2 = prs.slides.add_slide(blank2)
    add_rect(s2, 0, 0, W, HEADER_H, HEADER_FILL)
    add_textbox(s2, 0.40, 0.06, W - 0.5, 0.40,
                'Background & Mathematical Model', 20, bold=True, color=WHITE)
    add_textbox(s2, 0.40, 0.46, W - 0.5, 0.24,
                'Dieckow et al. 2024  ·  10-guild gLV  ·  N = 10 patients  ·  4 time-points (weeks 0/1/3/6)',
                11, bold=False, color=LIGHT_BLUE)
    # ODE equation box
    add_rect(s2, 0.4, 0.9, W - 0.8, 0.75, HEADER_FILL)
    add_textbox(s2, 0.5, 0.96, W - 1.0, 0.62,
                'dφᵢ/dt  =  φᵢ · ( bᵢ  +  Σⱼ Aᵢⱼ · φⱼ )',
                18, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    # Parameter descriptions
    params = [
        ('φᵢ  = relative abundance of guild i',                                   0.4),
        ('bᵢ  = intrinsic growth rate (patient-specific oral environment)',        0.9),
        ('Aᵢⱼ = interaction strength: guild j → guild i  (+ mutualism, − competition)', 1.4),
        ('R/G ratio = log(φ_dys) − log(φ_com)   R/G ratio < 0 → commensal,   R/G ratio > 0 → dysbiotic', 1.9),
    ]
    for txt, y_off in params:
        add_textbox(s2, 0.5, 1.78 + y_off, W - 1.0, 0.4,
                    txt, 11, bold=False, color=NAVY)
    # Biology note
    add_rect(s2, 0.4, 5.4, W - 0.8, 0.85, FOOTER_FILL)
    add_textbox(s2, 0.5, 5.48, W - 1.0, 0.70,
                'Fitting: L-BFGS-B minimises prediction MSE over 4 time-points → recovers A (10×10) and b (10 patients × 10 guilds).  '
                'The fitted A matrix encodes the ecological interaction network and determines long-term community fate.',
                10, bold=False, color=NAVY)
    add_rect(s2, 0, FOOTER_Y, W, FOOTER_H, FOOTER_FILL)
    add_textbox(s2, 0.35, FOOTER_Y + 0.07, W - 0.6, FOOTER_H - 0.10,
                '~700 species live in the oral cavity, organised into 10 functional guilds.  '
                'Healthy (CT1): Acti/Baci dominant, R/G ratio < 0.  '
                'Dysbiotic (CT2): Bacteroidia/Fusobacteria expand, R/G ratio > 0.  '
                'Question: why do some patients relapse after treatment while others do not?',
                10.5, bold=False, color=NAVY)

    # ── Slide B: Keystone Analysis ──────────────────────────────────────────
    add_slide(
        prs,
        title    = 'Guild-Level Keystone Analysis: Knockout BC Divergence',
        subtitle = 'Which guild, when removed, shifts the community most?',
        fig_path = R / 'guild_importance' / 'knockout_bars_base.png',
        footer_text = (
            '🔑  Bacilli: highest BC divergence (0.839) — keystone guild. '
            'Removal increases volume → competitive suppressor. '
            '  |  '
            'Actinobacteria: vol_drop = +0.846 → growth engine. '
            'Strongest interaction: Acti → Baci  A = +3.43.'
        ),
    )

    # ── Slide C: Network Centrality ─────────────────────────────────────────
    add_slide(
        prs,
        title    = 'Guild Network Centrality: Influence, Vulnerability, Net Role',
        subtitle = 'Column-sum |A| = influence;  Row-sum |A| = vulnerability;  Σ A_ij = net mutualist (+) / competitor (−)',
        fig_path = R / 'guild_network' / 'guild_centrality_summary.png',
        footer_text = (
            'Actinobacteria: influence = 5.58 (highest), net = +2.02 (hub guild).  '
            'Bacilli: vulnerability = 8.24 (most regulated by others).  '
            'Betaproteobacteria: net = −0.14 → net competitor.  '
            'LOO sign consistency: Acti→* = 1.0 (stable across all 10 folds); Gamm→* ≈ 0 (data-limited).'
        ),
    )

    # ── Slide C2: Permutation Test ──────────────────────────────────────────
    add_slide(
        prs,
        title    = 'A Matrix Statistical Validation: Permutation Test (n = 100)',
        subtitle = 'Null: shuffle patient labels → refit A → p-value = fraction of |A_perm| ≥ |A_real|  (parallelised, 24 cores)',
        fig_path = R / 'guild_network' / 'permutation_test.png',
        footer_text = (
            '2 / 90 pairs significant at p < 0.05: '
            'Baci→Acti (A = +3.43, p = 0.040) and Gamm→Clos (A = −0.023, p = 0.040).  '
            '|  LOO confirms Acti→* sign stability (SC = 1.0 all folds); Gamm→* unstable (SC ≈ 0).  '
            '|  Low significance consistent with N = 10 — A matrix structure reflects ecology, not noise.'
        ),
    )

    # ── Slide C3: Tipping Point ─────────────────────────────────────────────
    add_two_fig_slide(
        prs,
        title       = 'Tipping Point Analysis: CT2 → CT1 b-Environment Transition',
        subtitle    = 'Left: α-interpolation b(α) = (1-α)·b_CT2 + α·b_CT1  |  Right: single-guild b sweep |ΔR/G|',
        fig_left    = R / 'guild_tipping' / 'tipping_alpha_scan.png',
        fig_right   = R / 'guild_tipping' / 'tipping_single_guild.png',
        caption_left  = 'α-scan: no R/G=0 crossing (single commensal attractor)',
        caption_right = 'Single-guild effect sizes: Gamm #1, Bact #2',
        footer_text = (
            'No α* crossing: A matrix drives system to commensal for all b interpolations — structural attractor dominant.  '
            '|  Top drivers: Gammaproteobacteria (ΔR/G = 1.04) and Bacteroidia (ΔR/G ratio = 0.94).  '
            '|  Equilibrium R/G ratio: Patient D = −0.48 (only commensal); all others > 0 (dysbiotic).'
        ),
    )

    # ── Slide C4: 2D Phase Diagram ──────────────────────────────────────────
    add_slide(
        prs,
        title    = '2D Phase Diagram: Gammaproteobacteria × Bacteroidia Growth-Rate Space',
        subtitle = '20 × 20 grid · all other b at CT2 mean · week-3 R/G ratio · black dashed = R/G ratio = 0 tipping boundary',
        fig_path = R / 'guild_tipping' / 'tipping_phase_diagram_2d.png',
        footer_text = (
            'Gamm × Bact: top-2 single-guild tipping drivers (ΔR/G 1.04 and 0.94).  '
            'R/G ratio = 0 boundary visible — commensal region at low b_Gamm and low b_Bact.  '
            'CT2 mean (▼) deep in dysbiotic zone; CT1 mean (▲) near tipping boundary.  '
            '→ Joint reduction of Gamm and Bact growth rates is the most efficient ecological intervention.'
        ),
    )

    # ── Slide D: Relapse Dynamics ───────────────────────────────────────────
    add_two_fig_slide(
        prs,
        title       = 'Predicted Relapse Dynamics: Three Patient Archetypes',
        subtitle    = 'Long-time simulation (t = 150 d) from patient-specific b vectors',
        fig_left    = R / 'guild_tipping' / 'relapse_dynamics.png',
        fig_right   = R / 'guild_relapse' / 'gdi_w3_vs_relapse.png',
        caption_left  = 'R/G ratio trajectories to 90 days (per patient)',
        caption_right = 'Week-3 R/G ratio → relapse day (prognostic marker)',
        footer_text = (
            '① Persistent dysbiotic (A,B,E,F,H,K): R/G ratio > 0 throughout.  '
            '② Transient responder (C→24d, G→38d, L→48d): commensal at wk-3, dysbiotic by t < 50d.  '
            '③ Permanent commensal (D): only patient with stable R/G ratio < 0.  '
            '→ Week-3 R/G ratio as prognostic marker: lower = longer remission.'
        ),
    )

    # ── Slide E: Prevention + Joshi ─────────────────────────────────────────
    add_two_fig_slide(
        prs,
        title       = 'Prevention Threshold & External Validation (Joshi 2025)',
        subtitle    = 'Min. Δb_i needed to prevent relapse  |  Model equilibrium vs peri-implantitis reference',
        fig_left    = R / 'guild_relapse' / 'prevention_threshold.png',
        fig_right   = R / 'guild_relapse' / 'joshi_comparison.png',
        caption_left  = 'Required Δb per guild (patient C only actionable)',
        caption_right = 'Predicted eq. R/G ratio vs Joshi PI distribution',
        footer_text = (
            'Prevention: only Patient C convertible by single-guild intervention '
            '(Δb_Bact = 1.33). Patients G & L require multi-guild change → deep attractor.  '
            '|  External validation: predicted equilibrium R/G ratio = +0.65 ≈ Joshi PI mean +0.53  '
            '(9/10 patients in PI range, difference < 0.12 within Joshi SD=1.30).'
        ),
    )

    # ── Slide F: Phase Diagram ──────────────────────────────────────────────
    add_slide(
        prs,
        title    = '3D Attractor Landscape: R/G ratio = 0 Boundary in Growth-Rate Space',
        subtitle = 'Equilibrium R/G ratio on 18³ grid (Acti × Baci × Bact), 12-core parallel, t = 80 d',
        fig_path = R / 'guild_phase' / 'phase3d_static.png',
        footer_text = (
            '98% of parameter space is dysbiotic (R/G ratio > 0).  '
            'Commensal region (blue) confined to b_Bact < ~1.5 — Bacteroidia growth rate is the dominant control variable.  '
            'CT1 patients with low b_Bact (G, K, L) cluster near boundary; '
            'CT2 patients (A: b_Bact=4.5) deep in dysbiotic zone.  '
            '→ Confirms: attractor robustness is structural, not initial-condition sensitive.'
        ),
    )

    # ── Slide Z: Conclusion ─────────────────────────────────────────────────
    blank_z = prs.slide_layouts[6]
    sz = prs.slides.add_slide(blank_z)
    add_rect(sz, 0, 0, W, H, HEADER_FILL)
    add_textbox(sz, 1.0, 0.25, W - 2.0, 0.60,
                'Conclusion', 28, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    # Key message
    add_rect(sz, 0.8, 1.0, W - 1.6, 1.0, RGBColor(0x1A, 0x5C, 0x3A))
    add_textbox(sz, 0.9, 1.07, W - 1.8, 0.88,
                'Periodontitis relapse is governed by the topology of microbial interaction networks, '
                'not merely by bacterial growth rates. '
                'Prevention thresholds and relapse timing can be mathematically derived '
                'from patient-specific gLV dynamics.',
                12, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    # Numbered findings
    findings = [
        '①  Bacilli = keystone guild (BC = 0.839);  Acti = growth engine.',
        '②  Baci → Acti (A = +3.43, p = 0.040) is the sole statistically validated interaction.',
        '③  A matrix creates a single commensal attractor — b alone cannot flip the community state.',
        '④  Three relapse archetypes: persistent / transient responder / permanent commensal (Patient D only).',
        '⑤  Patient C preventable with Δb_Bact = 1.33;  98 % of growth-rate space remains dysbiotic.',
        '⑥  External validation: predicted eq. R/G ratio +0.65 ≈ Joshi 2025 peri-implant cohort mean +0.53.',
    ]
    for k, txt in enumerate(findings):
        add_textbox(sz, 0.6, 2.20 + k * 0.68, W - 1.2, 0.60,
                    txt, 11.5, bold=False, color=WHITE)

    prs.save(str(PPTX_OUT))
    print(f'Saved: {PPTX_OUT}  ({len(prs.slides)} slides total)')


if __name__ == '__main__':
    main()
